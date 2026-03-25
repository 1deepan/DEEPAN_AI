import os
import PyPDF2
from pptx import Presentation
import docx
import openpyxl

def extract_text_from_pdf(file_path):
    text = ""
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        print(f"[Parser] PDF Error: {e}")
    return text.strip()

def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"[Parser] PPTX Error: {e}")
    return text.strip()

def extract_text_from_docx(file_path):
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"[Parser] DOCX Error: {e}")
    return text.strip()

def extract_text_from_xlsx(file_path):
    text = ""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet in wb.worksheets:
            text += f"--- Sheet: {sheet.title} ---\n"
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
    except Exception as e:
        print(f"[Parser] XLSX Error: {e}")
    return text.strip()

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == '.pdf':
        return extract_text_from_pdf(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_text_from_pptx(file_path)
    elif ext == '.docx':
        return extract_text_from_docx(file_path)
    elif ext == '.xlsx':
        return extract_text_from_xlsx(file_path)
    elif ext in ['.txt', '.md', '.py', '.c', '.cpp', '.java', '.js', '.html', '.css']:
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except:
            return ""
    return ""
